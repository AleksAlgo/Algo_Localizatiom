"""
Извлечение и обратная сборка переводимых сегментов из docx / xlsx / pptx / xliff.

Каждый handler возвращает list[Segment]:
    Segment = {
        "id":   str,            # уникальный идентификатор сегмента
        "text": str,            # исходный текст сегмента (то, что переводим)
        "ctx":  dict (opt.),    # контекст для LLM (тип, окружение)
    }

И умеет принять обновлённые тексты (по тому же id) и сохранить файл с сохранением
форматирования: для docx/xlsx/pptx — на уровне run/cell/text-frame,
для xliff — заполнение <target>.
"""
from __future__ import annotations

import re
import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree as ET


@dataclass
class Segment:
    id: str
    text: str
    ctx: dict = field(default_factory=dict)


# ────────────────────────────────────────────────────────────────────────
# DOCX
# ────────────────────────────────────────────────────────────────────────
class DocxHandler:
    ext = ".docx"

    def __init__(self, path: Path):
        from docx import Document
        self.path = Path(path)
        self.doc = Document(self.path)

    @staticmethod
    def _hl_runs(p_el):
        """Yield (hi, ri, run_element) for runs inside <w:hyperlink> elements."""
        from docx.oxml.ns import qn
        for hi, hl in enumerate(p_el.findall(qn("w:hyperlink"))):
            for ri, r_el in enumerate(hl.findall(qn("w:r"))):
                yield hi, ri, r_el

    def extract(self) -> list[Segment]:
        from docx.oxml.ns import qn
        segments: list[Segment] = []
        # параграфы (включая заголовки и списки)
        for pi, p in enumerate(self.doc.paragraphs):
            for ri, run in enumerate(p.runs):
                if run.text and run.text.strip():
                    segments.append(Segment(
                        id=f"p{pi}.r{ri}",
                        text=run.text,
                        ctx={"style": p.style.name if p.style else "Normal"},
                    ))
            # гиперссылки в параграфе
            for hi, ri, r_el in self._hl_runs(p._p):
                t_el = r_el.find(qn("w:t"))
                if t_el is not None and t_el.text and t_el.text.strip():
                    segments.append(Segment(
                        id=f"p{pi}.hl{hi}.r{ri}",
                        text=t_el.text,
                        ctx={"style": p.style.name if p.style else "Normal", "hyperlink": True},
                    ))
        # таблицы
        for ti, table in enumerate(self.doc.tables):
            for row_i, row in enumerate(table.rows):
                for col_i, cell in enumerate(row.cells):
                    for pi, p in enumerate(cell.paragraphs):
                        for ri, run in enumerate(p.runs):
                            if run.text and run.text.strip():
                                segments.append(Segment(
                                    id=f"t{ti}.r{row_i}.c{col_i}.p{pi}.r{ri}",
                                    text=run.text,
                                    ctx={"loc": "table"},
                                ))
                        # гиперссылки в ячейке
                        for hi, ri, r_el in self._hl_runs(p._p):
                            t_el = r_el.find(qn("w:t"))
                            if t_el is not None and t_el.text and t_el.text.strip():
                                segments.append(Segment(
                                    id=f"t{ti}.r{row_i}.c{col_i}.p{pi}.hl{hi}.r{ri}",
                                    text=t_el.text,
                                    ctx={"loc": "table", "hyperlink": True},
                                ))
        return segments

    def write(self, translated: dict[str, str], out_path: Path):
        from docx.oxml.ns import qn
        # параграфы
        for pi, p in enumerate(self.doc.paragraphs):
            for ri, run in enumerate(p.runs):
                key = f"p{pi}.r{ri}"
                if key in translated:
                    run.text = translated[key]
            # гиперссылки в параграфе
            for hi, ri, r_el in self._hl_runs(p._p):
                key = f"p{pi}.hl{hi}.r{ri}"
                if key in translated:
                    t_el = r_el.find(qn("w:t"))
                    if t_el is not None:
                        t_el.text = translated[key]
        # таблицы
        for ti, table in enumerate(self.doc.tables):
            for row_i, row in enumerate(table.rows):
                for col_i, cell in enumerate(row.cells):
                    for pi, p in enumerate(cell.paragraphs):
                        for ri, run in enumerate(p.runs):
                            key = f"t{ti}.r{row_i}.c{col_i}.p{pi}.r{ri}"
                            if key in translated:
                                run.text = translated[key]
                        # гиперссылки в ячейке
                        for hi, ri, r_el in self._hl_runs(p._p):
                            key = f"t{ti}.r{row_i}.c{col_i}.p{pi}.hl{hi}.r{ri}"
                            if key in translated:
                                t_el = r_el.find(qn("w:t"))
                                if t_el is not None:
                                    t_el.text = translated[key]
        self.doc.save(out_path)


# ────────────────────────────────────────────────────────────────────────
# XLSX
# ────────────────────────────────────────────────────────────────────────
class XlsxHandler:
    ext = ".xlsx"

    def __init__(self, path: Path):
        from openpyxl import load_workbook
        self.path = Path(path)
        self.wb = load_workbook(self.path)

    def extract(self) -> list[Segment]:
        segments: list[Segment] = []
        for ws_name in self.wb.sheetnames:
            ws = self.wb[ws_name]
            for row in ws.iter_rows():
                for cell in row:
                    val = cell.value
                    if isinstance(val, str) and val.strip() and not val.startswith("="):
                        segments.append(Segment(
                            id=f"{ws_name}!{cell.coordinate}",
                            text=val,
                            ctx={"sheet": ws_name, "coord": cell.coordinate},
                        ))
        return segments

    def write(self, translated: dict[str, str], out_path: Path):
        for key, new_text in translated.items():
            sheet, coord = key.split("!", 1)
            self.wb[sheet][coord] = new_text
        self.wb.save(out_path)


# ────────────────────────────────────────────────────────────────────────
# PPTX
# ────────────────────────────────────────────────────────────────────────
class PptxHandler:
    ext = ".pptx"

    def __init__(self, path: Path):
        from pptx import Presentation
        self.path = Path(path)
        self.prs = Presentation(self.path)

    @staticmethod
    def _extract_tf(tf, prefix: str, ctx: dict) -> list[Segment]:
        segs = []
        for pi, para in enumerate(tf.paragraphs):
            for ri, run in enumerate(para.runs):
                if run.text and run.text.strip():
                    segs.append(Segment(
                        id=f"{prefix}.p{pi}.r{ri}",
                        text=run.text,
                        ctx=ctx,
                    ))
        return segs

    @staticmethod
    def _write_tf(tf, prefix: str, translated: dict[str, str]):
        for pi, para in enumerate(tf.paragraphs):
            for ri, run in enumerate(para.runs):
                key = f"{prefix}.p{pi}.r{ri}"
                if key in translated:
                    run.text = translated[key]

    def extract(self) -> list[Segment]:
        segments: list[Segment] = []
        for si, slide in enumerate(self.prs.slides):
            for shi, shape in enumerate(slide.shapes):
                prefix = f"s{si}.sh{shi}"
                if shape.has_text_frame:
                    segments.extend(self._extract_tf(shape.text_frame, prefix, {"slide": si + 1}))
                elif shape.shape_type == 19:  # TABLE
                    for row_i, row in enumerate(shape.table.rows):
                        for col_i, cell in enumerate(row.cells):
                            tbl_prefix = f"{prefix}.tbl.row{row_i}.col{col_i}"
                            segments.extend(self._extract_tf(cell.text_frame, tbl_prefix, {"slide": si + 1, "loc": "table"}))
        return segments

    def write(self, translated: dict[str, str], out_path: Path):
        for si, slide in enumerate(self.prs.slides):
            for shi, shape in enumerate(slide.shapes):
                prefix = f"s{si}.sh{shi}"
                if shape.has_text_frame:
                    self._write_tf(shape.text_frame, prefix, translated)
                elif shape.shape_type == 19:  # TABLE
                    for row_i, row in enumerate(shape.table.rows):
                        for col_i, cell in enumerate(row.cells):
                            tbl_prefix = f"{prefix}.tbl.row{row_i}.col{col_i}"
                            self._write_tf(cell.text_frame, tbl_prefix, translated)
        self.prs.save(out_path)


# ────────────────────────────────────────────────────────────────────────
# XLIFF (1.2 / 2.0)
# ────────────────────────────────────────────────────────────────────────
class XliffHandler:
    ext = ".xliff"
    NS = {
        "x12": "urn:oasis:names:tc:xliff:document:1.2",
        "x20": "urn:oasis:names:tc:xliff:document:2.0",
    }

    def __init__(self, path: Path):
        self.path = Path(path)
        self.tree = ET.parse(self.path)
        self.root = self.tree.getroot()
        self.version = "2.0" if self.root.tag.endswith("}xliff") and "2.0" in self.root.tag else "1.2"

    def _iter_units(self):
        if "2.0" in self.root.tag:
            for unit in self.root.iter("{urn:oasis:names:tc:xliff:document:2.0}unit"):
                for seg in unit.iter("{urn:oasis:names:tc:xliff:document:2.0}segment"):
                    yield seg, "x20"
        else:
            for tu in self.root.iter("{urn:oasis:names:tc:xliff:document:1.2}trans-unit"):
                yield tu, "x12"

    def extract(self) -> list[Segment]:
        segments: list[Segment] = []
        for i, (unit, ns_key) in enumerate(self._iter_units()):
            ns = self.NS[ns_key]
            src_el = unit.find(f"{{{ns}}}source")
            if src_el is None:
                continue
            src_text = "".join(src_el.itertext())
            uid = unit.get("id") or f"unit{i}"
            segments.append(Segment(id=uid, text=src_text, ctx={"xliff_version": self.version}))
        return segments

    def write(self, translated: dict[str, str], out_path: Path):
        for i, (unit, ns_key) in enumerate(self._iter_units()):
            ns = self.NS[ns_key]
            uid = unit.get("id") or f"unit{i}"
            if uid not in translated:
                continue
            tgt_el = unit.find(f"{{{ns}}}target")
            if tgt_el is None:
                tgt_el = ET.SubElement(unit, f"{{{ns}}}target")
            tgt_el.text = translated[uid]
        # сохраняем с XML declaration
        self.tree.write(out_path, encoding="utf-8", xml_declaration=True)


# ────────────────────────────────────────────────────────────────────────
# TXT / MD (бонус — простые форматы)
# ────────────────────────────────────────────────────────────────────────
class TextHandler:
    ext = ".txt"

    def __init__(self, path: Path):
        self.path = Path(path)
        self.lines = self.path.read_text(encoding="utf-8").splitlines(keepends=True)

    def extract(self) -> list[Segment]:
        segs = []
        for i, line in enumerate(self.lines):
            stripped = line.strip()
            if stripped:
                segs.append(Segment(id=f"L{i}", text=stripped, ctx={"line": i}))
        return segs

    def write(self, translated: dict[str, str], out_path: Path):
        out_lines = list(self.lines)
        for i, line in enumerate(self.lines):
            key = f"L{i}"
            if key in translated:
                # сохраняем ведущие/завершающие пробелы и переносы
                lead = re.match(r"^(\s*)", line).group(1)
                trail = re.search(r"(\s*)$", line).group(1)
                out_lines[i] = f"{lead}{translated[key]}{trail}"
        Path(out_path).write_text("".join(out_lines), encoding="utf-8")


# ────────────────────────────────────────────────────────────────────────
# Фабрика
# ────────────────────────────────────────────────────────────────────────
HANDLERS = {
    ".docx":  DocxHandler,
    ".xlsx":  XlsxHandler,
    ".pptx":  PptxHandler,
    ".xliff": XliffHandler,
    ".xlf":   XliffHandler,
    ".txt":   TextHandler,
    ".md":    TextHandler,
}


def get_handler(path: Path):
    ext = Path(path).suffix.lower()
    if ext not in HANDLERS:
        raise ValueError(f"Unsupported file format: {ext}. Supported: {list(HANDLERS)}")
    return HANDLERS[ext](path)
